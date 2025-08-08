import streamlit as st
from dotenv import load_dotenv
import requests
from bs4 import BeautifulSoup
import plotly.graph_objects as go
from ai_analyzer import analyze_competency_gemini, ocr_with_gemini
import PyPDF2
from pptx import Presentation
import io
import os
from streamlit.errors import StreamlitSecretNotFoundError
from streamlit_paste_button import paste_image_button

# --- 페이지 설정 ---
st.set_page_config(layout="wide", page_title="JOJUN - AI 직무 역량 조준기")

# --- 스타일링 ---
st.markdown("""
<style>
    /* 폰트 */
    @import url('https://fonts.googleapis.com/css2?family=Pretendard:wght@400;600;700&display=swap');
    .stApp {
        font-family: 'Pretendard', sans-serif;
    }

    /* === 버튼 스타일 === */
    .stButton>button {
        font-family: 'Pretendard', sans-serif;
        font-weight: 700;
        font-size: 16px;
        color: white;
        background-color: #4A4A4A;
        border: none;
        border-radius: 10px;
        padding: 12px 0;
        transition: all 0.2s ease-in-out;
    }
    .stButton>button:hover {
        background-color: #2a2a2a;
        transform: scale(1.02);
    }
    .stButton>button:active {
        background-color: #1a1a1a !important;
        transform: scale(0.98) !important;
        color: white !important;
    }

    /* === 라이트/다크 모드 스타일 === */
    .kpi-card { background-color: #FFFFFF; border-left: 5px solid #4A4A4A; color: #333; }
    .kpi-title { color: #333333; }
    .kpi-score-label { color: #666666; }
    .kpi-score-value { color: #333; }
    .ai-comment-card { background-color: #f8f9fa; border: 1px solid #dee2e6; }
    .ai-comment-title { color: #4A4A4A; }
    .ai-comment-body { color: #343a40; }
    [data-theme="dark"] .kpi-card { background-color: #262730; border-left: 5px solid #999; color: #FAFAFA; }
    [data-theme="dark"] .kpi-title { color: #FAFAFA; }
    [data-theme="dark"] .kpi-score-label { color: #A0A0A0; }
    [data-theme="dark"] .kpi-score-value { color: #FAFAFA; }
    [data-theme="dark"] .ai-comment-card { background-color: #262730; border: 1px solid #444; }
    [data-theme="dark"] .ai-comment-title { color: #FAFAFA; }
    [data-theme="dark"] .ai-comment-body { color: #A0A0A0; }

    /* === 공통 스타일 === */
    .kpi-card { border-radius: 10px; padding: 20px; box-shadow: 0 4px 8px rgba(0,0,0,0.1); margin-bottom: 10px; height: 180px; display: flex; flex-direction: column; justify-content: space-between; }
    .kpi-title { font-size: 1.1rem; font-weight: 700; margin-bottom: 15px; }
    .kpi-scores { display: flex; justify-content: space-between; align-items: center; }
    .kpi-score-box { text-align: center; }
    .kpi-score-label { font-size: 0.8rem; }
    .kpi-score-value { font-size: 1.8rem; font-weight: 700; }
    .kpi-delta { text-align: center; font-size: 1.2rem; font-weight: 700; }
    .delta-positive { color: #28a745 !important; }
    .delta-negative { color: #dc3545 !important; }
    .delta-zero { color: #6c757d !important; }
    .ai-comment-card { border-radius: 10px; padding: 25px; margin-top: 5px; }
    .ai-comment-title { font-size: 1.2rem; font-weight: 700; margin-bottom: 10px; }
    .ai-comment-body { font-size: 1rem; line-height: 1.6; }
</style>
""", unsafe_allow_html=True)

# --- 함수 정의 ---
def parse_input_files(uploaded_files):
    all_text = ""
    if not uploaded_files:
        return all_text
    
    progress_bar = st.sidebar.progress(0)
    for i, file in enumerate(uploaded_files):
        try:
            file_extension = file.name.split('.')[-1].lower()
            file_bytes = file.getvalue()

            if file_extension == "pdf":
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
                for page in pdf_reader.pages:
                    all_text += page.extract_text() + "\n"
            elif file_extension == "pptx":
                presentation = Presentation(io.BytesIO(file_bytes))
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            all_text += shape.text + "\n"
            elif file_extension in ["jpg", "jpeg", "png"]:
                with st.spinner(f"'{file.name}' 이미지 텍스트 분석 중..."):
                    ocr_text = ocr_with_gemini(file_bytes)
                if ocr_text:
                    all_text += ocr_text + "\n"
            else:
                all_text += file_bytes.decode("utf-8", errors='ignore') + "\n"
        except Exception as e:
            st.sidebar.error(f"'{file.name}' 처리 중 오류: {e}")
        finally:
            progress_bar.progress((i + 1) / len(uploaded_files), f"'{file.name}' 처리 완료!")
    progress_bar.empty()
    st.sidebar.success(f"{len(uploaded_files)}개 파일 분석 완료!")
    return all_text

# --- session_state 초기화 ---
if 'last_pasted_image' not in st.session_state:
    st.session_state.last_pasted_image = None
if 'jd_text' not in st.session_state:
    st.session_state.jd_text = ""


# --- UI - 사이드바 (입력) ---
with st.sidebar:
    st.title("🎯 JOJUN")
    st.markdown("AI 직무 역량 조준기")
    
    st.divider()

    st.header("1. 채용 공고 입력")
    
    with st.expander("🔗 URL에서 가져오기", expanded=False):
        url_input = st.text_input("채용 공고 URL")
        if url_input:
            try:
                headers = {'User-Agent': 'Mozilla/5.0'}
                response = requests.get(url_input, headers=headers, timeout=5)
                response.raise_for_status()
                soup = BeautifulSoup(response.content, 'html.parser')
                st.session_state.jd_text = soup.find('body').get_text(separator='\n', strip=True)
                st.success("URL 내용을 성공적으로 가져왔습니다.")
            except Exception as e:
                st.error(f"URL 처리 오류: {e}")

    with st.expander("✍️ 붙여넣기 & 직접 수정", expanded=True):
        paste_result = paste_image_button("📋 클립보드 이미지 붙여넣기", key="paste_button")
        
        # === 로직 수정: 새 이미지가 붙여넣어졌을 때만 OCR 실행 ===
        if paste_result.image_data and paste_result.image_data != st.session_state.last_pasted_image:
            st.session_state.last_pasted_image = paste_result.image_data
            with st.spinner("이미지 분석 중..."):
                img_bytes_io = io.BytesIO()
                paste_result.image_data.save(img_bytes_io, format="PNG")
                ocr_text = ocr_with_gemini(img_bytes_io.getvalue())

            if ocr_text:
                st.session_state.jd_text += "\n" + ocr_text
                st.info("이미지 텍스트를 공고 내용에 추가했습니다.")
                st.rerun() # UI 즉시 새로고침

        st.text_area("공고 내용", key="jd_text", height=200)

    with st.expander("📁 파일 업로드", expanded=False):
        jd_files = st.file_uploader(
            "PDF, PPTX, TXT, MD, JPG, PNG",
            type=["pdf", "pptx", "txt", "md", "jpg", "jpeg", "png"],
            accept_multiple_files=True,
            key="jd_uploader"
        )
    
    st.divider()

    st.header("2. 나의 경험 입력")
    with st.expander("✍️ 직접 입력", expanded=True):
        st.text_area("이력서, 포트폴리오 내용", key="my_exp_text", height=200)
    
    with st.expander("📁 파일 업로드", expanded=False):
        my_files = st.file_uploader(
            "PDF, PPTX, TXT, MD",
            type=["pdf", "pptx", "txt", "md"],
            accept_multiple_files=True,
            key="my_files_uploader"
        )

    st.divider()
    
    analyze_button = st.button("✨ AI로 합격률 조준하기", use_container_width=True)


# --- 메인 화면 (결과) ---
# (이하 코드 수정 없음)
st.title("🎯 JOJUN: AI 직무 역량 분석 결과")
st.markdown("채용 공고와 당신의 경험을 분석하여 합격 가능성을 알려드립니다.")

if 'analysis_data' not in st.session_state:
    st.session_state.analysis_data = None

if analyze_button:
    try:
        _ = st.secrets["GOOGLE_API_KEY"]
    except (StreamlitSecretNotFoundError, KeyError):
        if "GOOGLE_API_KEY" not in os.environ:
            st.error("Google API 키가 설정되지 않았습니다. .env 파일 또는 Streamlit Secrets에 추가해주세요.")
            st.stop()
            
    final_jd_text = st.session_state.get('jd_text', '')
    if jd_files:
        final_jd_text += "\n" + parse_input_files(jd_files)

    final_my_exp_text = st.session_state.get('my_exp_text', '')
    if my_files:
        final_my_exp_text += "\n" + parse_input_files(my_files)

    if not final_jd_text.strip() or not final_my_exp_text.strip():
        st.warning("채용 공고와 나의 경험을 모두 입력(또는 업로드)해주세요.")
        st.session_state.analysis_data = None
    else:
        with st.spinner("AI 'JOJUN'이 당신의 역량을 정밀 분석 중입니다..."):
            st.session_state.analysis_data = analyze_competency_gemini(final_jd_text, final_my_exp_text)

if st.session_state.analysis_data:
    analysis_data = st.session_state.analysis_data
    st.success("🎉 분석 완료! 당신의 합격 포인트를 확인하세요.")
    
    st.subheader("📊 종합 분석")
    fit_score = int(analysis_data.get('fit_score', 0))
    comment = analysis_data.get('overall_comment', "코멘트를 생성하지 못했습니다.")
    
    col1, col2 = st.columns([1, 2])
    with col1:
        st.metric(label="직무 적합도", value=f"{fit_score}점")
        st.progress(fit_score)
    with col2:
        st.markdown(f"""
        <div class="ai-comment-card">
            <div class="ai-comment-title">💡 AI 총평</div>
            <div class="ai-comment-body">{comment}</div>
        </div>
        """, unsafe_allow_html=True)
    
    st.divider()

    st.subheader("📈 역량 비교 분석")
    categories = analysis_data['categories']
    job_scores = analysis_data['job_scores']
    user_scores = analysis_data['user_scores']

    fig = go.Figure()
    fig.add_trace(go.Scatterpolar(r=job_scores, theta=categories, fill='toself', name='요구 역량 (JD)', line_color='rgba(74, 74, 74, 0.8)', fillcolor='rgba(74, 74, 74, 0.2)'))
    fig.add_trace(go.Scatterpolar(r=user_scores, theta=categories, fill='toself', name='보유 역량 (나)', line_color='rgba(255, 140, 0, 0.8)', fillcolor='rgba(255, 140, 0, 0.2)'))
    fig.update_layout(polar=dict(radialaxis=dict(visible=True, range=[0, 100], showline=False, showticklabels=False, ticks='')), showlegend=True, title=dict(text="<b>역량 적합도 레이더 차트</b>", font=dict(size=20), x=0.5), font=dict(family="Pretendard, sans-serif", size=14), legend=dict(yanchor="top", y=1.1, xanchor="center", x=0.5, orientation="h"), template="plotly_white", margin=dict(t=80, b=20))
    st.plotly_chart(fig, use_container_width=True)
    
    st.subheader("🔍 역량별 상세 점수")
    num_categories = len(categories)
    cols = st.columns(min(num_categories, 3))
    
    for i, category in enumerate(categories):
        job_score = job_scores[i]
        user_score = user_scores[i]
        delta = user_score - job_score
        
        if delta > 0: delta_color, delta_sign = "positive", "+"
        elif delta < 0: delta_color, delta_sign = "negative", ""
        else: delta_color, delta_sign = "zero", ""

        with cols[i % min(num_categories, 3)]:
            st.markdown(f"""
            <div class="kpi-card">
                <div class="kpi-title">{category}</div>
                <div class="kpi-scores">
                    <div class="kpi-score-box"><div class="kpi-score-label">요구 역량</div><div class="kpi-score-value">{job_score}</div></div>
                    <div class="kpi-score-box"><div class="kpi-score-label">보유 역량</div><div class="kpi-score-value">{user_score}</div></div>
                    <div class="kpi-delta"><div class="kpi-score-label">차이</div><div class="{delta_color}">{delta_sign}{delta}</div></div>
                </div>
            </div>
            """, unsafe_allow_html=True)
else:
    if not analyze_button:
        st.info("사이드바에서 채용 공고와 자신의 경험을 입력한 후, 'AI로 합격률 조준하기' 버튼을 눌러주세요.")